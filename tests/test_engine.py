"""Tests for the V2 Rendering Engine — elements, layout, and renderer."""

import os
import tempfile
import pytest
from pathlib import Path

from slidecraft.engine.elements import (
    ElementType, Position, TextStyle, ShapeStyle,
    GradientSpec, GradientStop,
    SlideElement, SlideSpec, PresentationSpec,
    title_element, subtitle_element, text_element,
    bullets_element, shape_element, image_element, divider_element,
)
from slidecraft.engine.layout import (
    auto_fit_title_size,
    auto_fit_bullet_size,
    auto_fit_spacing,
    compute_two_column_layout,
    compute_timeline_positions,
    ContentMetrics,
    measure_content,
)
from slidecraft.engine.renderer import (
    SlideRenderer, render_presentation,
    hex_to_rgb, hex_to_srgb,
    apply_gradient_background,
    apply_shadow,
)


# ══════════════════════════════════════════════════════════════
#  Element DSL Tests
# ══════════════════════════════════════════════════════════════

class TestElementDSL:
    """Test the declarative element builders."""

    def test_title_element(self):
        el = title_element("Hello World", left=1, top=2)
        assert el.element_type == ElementType.TITLE
        assert el.content == "Hello World"
        assert el.position.left == 1
        assert el.position.top == 2
        assert el.text_style.bold is True

    def test_subtitle_element(self):
        el = subtitle_element("Subtitle", font_size=24)
        assert el.element_type == ElementType.SUBTITLE
        assert el.text_style.font_size == 24

    def test_bullets_element(self):
        items = ["Point 1", "Point 2", "Point 3"]
        el = bullets_element(items, left=1, top=2, width=10, height=5)
        assert el.element_type == ElementType.BULLETS
        assert len(el.items) == 3
        assert el.text_style.auto_fit is True

    def test_shape_element_with_shadow(self):
        el = shape_element("rounded_rect", 0, 0, 5, 3, fill_color="#336699", shadow=True)
        assert el.shape_style.shadow is True
        assert el.shape_style.fill_color == "#336699"

    def test_divider_element(self):
        el = divider_element(left=1, top=2, width=5, color="#ff0000")
        assert el.element_type == ElementType.DIVIDER
        assert el.shape_style.fill_color == "#ff0000"

    def test_position_shift(self):
        pos = Position(left=1, top=2, width=5, height=3)
        shifted = pos.shift(dx=2, dy=-1)
        assert shifted.left == 3
        assert shifted.top == 1
        assert shifted.width == 5  # unchanged

    def test_slide_spec_fluent_api(self):
        spec = SlideSpec(slide_number=1, background_color="#0f172a")
        spec.add(title_element("Test")).add(shape_element("rectangle"))
        assert len(spec.elements) == 2


class TestGradientSpec:
    """Test gradient specification."""

    def test_gradient_creation(self):
        g = GradientSpec(
            angle=225,
            stops=[
                GradientStop(color="#0f172a", position=0.0),
                GradientStop(color="#1e293b", position=0.5),
                GradientStop(color="#312e81", position=1.0),
            ],
        )
        assert g.angle == 225
        assert len(g.stops) == 3
        assert g.stops[0].position == 0.0
        assert g.stops[2].position == 1.0


# ══════════════════════════════════════════════════════════════
#  Layout Engine Tests
# ══════════════════════════════════════════════════════════════

class TestAutoFitTitle:
    """Test auto-fit title sizing."""

    def test_short_title_gets_bigger(self):
        size = auto_fit_title_size("标题", 36)
        assert size > 36  # Short titles get larger

    def test_normal_title_stays(self):
        size = auto_fit_title_size("这是一个正常长度的标题", 36)
        assert 32 <= size <= 38

    def test_long_title_shrinks(self):
        size = auto_fit_title_size("这是一个非常非常非常长的标题，包含很多很多的文字信息需要进行缩放处理", 36)
        assert size < 36

    def test_minimum_size(self):
        size = auto_fit_title_size("A" * 200, 36)
        assert size >= 20


class TestAutoFitBullets:
    """Test auto-fit bullet sizing."""

    def test_few_bullets_normal_size(self):
        items = ["P1", "P2", "P3"]
        size = auto_fit_bullet_size(items, 18, 5.0)
        assert size == 18  # Few items → no reduction

    def test_many_bullets_shrink(self):
        items = [f"Point {i}" for i in range(12)]
        size = auto_fit_bullet_size(items, 18, 5.0)
        assert size < 18  # Many items → shrink to fit

    def test_minimum_bullet_size(self):
        items = [f"P{i}" for i in range(30)]
        size = auto_fit_bullet_size(items, 18, 3.0)
        assert size >= 12  # Never below 12pt

    def test_long_bullets_reduce(self):
        items = ["A" * 90]  # Very long bullet
        size = auto_fit_bullet_size(items, 18, 5.0)
        assert size <= 16

    def test_empty_list(self):
        size = auto_fit_bullet_size([], 18, 5.0)
        assert size == 18


class TestAutoFitSpacing:
    """Test auto-fit bullet spacing."""

    def test_few_items_generous_spacing(self):
        spacing = auto_fit_spacing(2, 5.0, 18)
        assert spacing >= 12

    def test_many_items_tight_spacing(self):
        spacing = auto_fit_spacing(15, 5.0, 18)
        assert spacing <= 12

    def test_spacing_bounds(self):
        spacing = auto_fit_spacing(1, 5.0, 18)
        assert 4 <= spacing <= 24


class TestTwoColumnLayout:
    """Test smart two-column layout computation."""

    def test_balanced_columns(self):
        layout = compute_two_column_layout(3, 3)
        left_w = layout["left_card"].width
        right_w = layout["right_card"].width
        assert abs(left_w - right_w) < 0.1  # Equal items → equal width

    def test_unbalanced_favors_more_content(self):
        layout = compute_two_column_layout(5, 2)
        assert layout["left_card"].width > layout["right_card"].width

    def test_clamped_ratio(self):
        layout = compute_two_column_layout(10, 1)
        ratio = layout["left_card"].width / (layout["left_card"].width + layout["right_card"].width + 0.4)
        assert ratio <= 0.65  # Should be clamped


class TestTimelinePositions:
    """Test timeline milestone positioning."""

    def test_single_milestone_centered(self):
        positions = compute_timeline_positions(1)
        assert len(positions) == 1

    def test_multiple_milestones_spacing(self):
        positions = compute_timeline_positions(5)
        assert len(positions) == 5
        # Each milestone should be further right than the previous
        for i in range(1, 5):
            assert positions[i]["cx"] > positions[i - 1]["cx"]

    def test_alternating_above_below(self):
        positions = compute_timeline_positions(4)
        assert positions[0]["is_above"] is True
        assert positions[1]["is_above"] is False
        assert positions[2]["is_above"] is True
        assert positions[3]["is_above"] is False


class TestContentMetrics:
    """Test content analysis."""

    def test_empty_spec(self):
        spec = SlideSpec()
        m = measure_content(spec)
        assert m.content_density == 0.0

    def test_dense_content(self):
        spec = SlideSpec()
        spec.add(title_element("Test"))
        spec.add(bullets_element(["A" * 80 for _ in range(8)]))
        m = measure_content(spec)
        assert m.bullet_count == 8
        assert m.content_density > 0.3


# ══════════════════════════════════════════════════════════════
#  Renderer Tests
# ══════════════════════════════════════════════════════════════

class TestRendererUtils:
    """Test renderer utility functions."""

    def test_hex_to_rgb(self):
        from pptx.dml.color import RGBColor
        color = hex_to_rgb("#6366f1")
        assert color == RGBColor(0x63, 0x66, 0xf1)

    def test_hex_to_rgb_shorthand(self):
        from pptx.dml.color import RGBColor
        color = hex_to_rgb("#fff")
        assert color == RGBColor(0xff, 0xff, 0xff)

    def test_hex_to_srgb_strips_hash(self):
        assert hex_to_srgb("#6366f1") == "6366F1"


class TestSlideRendering:
    """Integration tests for full slide rendering."""

    def test_render_simple_slide(self):
        """Render a basic content slide."""
        spec = PresentationSpec(title="Test")
        slide_spec = SlideSpec(
            slide_number=1,
            background_color="#0f172a",
        )
        slide_spec.add(title_element("Hello", color="#ffffff"))
        slide_spec.add(bullets_element(["Point 1", "Point 2"], color="#cccccc"))
        spec.slides.append(slide_spec)

        prs = render_presentation(spec, {"text_primary": "#ffffff", "text_secondary": "#999999"})
        assert len(prs.slides) == 1

    def test_render_with_gradient(self):
        """Render a slide with gradient background."""
        spec = PresentationSpec(title="Gradient Test")
        slide_spec = SlideSpec(
            slide_number=1,
            background_gradient=GradientSpec(
                angle=225,
                stops=[
                    GradientStop(color="#0f172a", position=0.0),
                    GradientStop(color="#312e81", position=1.0),
                ],
            ),
        )
        slide_spec.add(title_element("Gradient!", color="#ffffff"))
        spec.slides.append(slide_spec)

        prs = render_presentation(spec, {"text_primary": "#ffffff"})
        assert len(prs.slides) == 1

    def test_render_with_shapes(self):
        """Render shapes with shadow effects."""
        spec = PresentationSpec(title="Shapes")
        slide_spec = SlideSpec(slide_number=1, background_color="#0f172a")
        slide_spec.add(shape_element("rounded_rect", 1, 1, 5, 3,
                                     fill_color="#1e293b", shadow=True))
        slide_spec.add(shape_element("oval", 7, 1, 2, 2,
                                     fill_color="#6366f1"))
        spec.slides.append(slide_spec)

        prs = render_presentation(spec, {})
        slide = prs.slides[0]
        assert len(slide.shapes) >= 2

    def test_render_z_order(self):
        """Test that elements are rendered in z_order."""
        spec = PresentationSpec(title="Z-Order")
        slide_spec = SlideSpec(slide_number=1)
        slide_spec.add(SlideElement(
            element_type=ElementType.SHAPE,
            position=Position(left=0, top=0, width=13, height=7),
            shape_style=ShapeStyle(fill_color="#000000"),
            z_order=-1,  # Background
        ))
        el = title_element("Front", color="#ffffff")
        el.z_order = 10
        slide_spec.add(el)  # Foreground
        spec.slides.append(slide_spec)

        prs = render_presentation(spec, {"text_primary": "#ffffff"})
        assert len(prs.slides) == 1

    def test_render_notes(self):
        """Test that speaker notes are preserved."""
        spec = PresentationSpec(title="Notes")
        slide_spec = SlideSpec(slide_number=1, notes="These are my notes")
        slide_spec.add(title_element("Title", color="#000"))
        spec.slides.append(slide_spec)

        prs = render_presentation(spec, {"text_primary": "#000"})
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "These are my notes" in notes

    def test_full_presentation_save(self, tmp_path):
        """Test saving a complete presentation to disk."""
        spec = PresentationSpec(title="Full Test")
        for i in range(5):
            ss = SlideSpec(slide_number=i + 1, background_color="#0f172a")
            ss.add(title_element(f"Slide {i + 1}", color="#fff"))
            spec.slides.append(ss)

        prs = render_presentation(spec, {"text_primary": "#fff"})
        out = tmp_path / "test_v2.pptx"
        prs.save(str(out))
        assert out.exists()
        assert out.stat().st_size > 0


class TestBuilderV2Integration:
    """Integration test for BuilderV2Agent."""

    def test_build_all_layouts(self, tmp_path):
        """Build a presentation with all 12 layout types using V2."""
        from slidecraft.agents.builder_v2 import BuilderV2Agent
        from slidecraft.config import SlideCraftConfig
        from slidecraft.models import (
            DesignSpec, DesignStyle, PresentationPlan,
            PresentationContent, SlideOutline, SlideContent, SlideType,
        )

        config = SlideCraftConfig()
        builder = BuilderV2Agent(config)
        design = DesignSpec(style=DesignStyle.DARK, use_gradient_backgrounds=True)

        plan = PresentationPlan(
            title="V2 Test", subtitle="Testing all layouts",
            slides=[
                SlideOutline(slide_number=i + 1, slide_type=st, title=f"S{i}")
                for i, st in enumerate([
                    SlideType.COVER, SlideType.TABLE_OF_CONTENTS,
                    SlideType.SECTION_HEADER, SlideType.CONTENT,
                    SlideType.TWO_COLUMN, SlideType.IMAGE_FULL,
                    SlideType.IMAGE_TEXT, SlideType.CHART,
                    SlideType.COMPARISON, SlideType.TIMELINE,
                    SlideType.QUOTE, SlideType.THANK_YOU,
                ])
            ],
        )

        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="V2 Test", subtitle="Enhanced builder"),
            SlideContent(slide_number=2, title="Contents"),
            SlideContent(slide_number=3, title="Section", subtitle="Part 01"),
            SlideContent(slide_number=4, title="Content", bullet_points=["A", "B", "C"]),
            SlideContent(slide_number=5, title="Columns", bullet_points=["L1", "L2", "R1", "R2"]),
            SlideContent(slide_number=6, title="Full Image"),
            SlideContent(slide_number=7, title="Image+Text", bullet_points=["P1", "P2"]),
            SlideContent(slide_number=8, title="Chart", data={
                "labels": ["A", "B", "C"], "values": [10, 20, 30],
            }),
            SlideContent(slide_number=9, title="Compare", bullet_points=["L", "R"]),
            SlideContent(slide_number=10, title="Timeline", bullet_points=["M1", "M2", "M3"]),
            SlideContent(slide_number=11, title="Quote", body_text="Test quote", subtitle="Author"),
            SlideContent(slide_number=12, title="Thanks!", subtitle="End"),
        ])

        out = str(tmp_path / "v2_all_layouts.pptx")
        result = builder.build(plan, content, design, [], out)

        assert Path(result).exists()
        from pptx import Presentation
        prs = Presentation(result)
        assert len(prs.slides) == 12
