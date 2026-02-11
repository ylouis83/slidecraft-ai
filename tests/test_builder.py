"""
Tests for the Builder Agent — verifies .pptx generation for all slide types.
"""

import os
import tempfile
import pytest

from pptx import Presentation as PptxPresentation

from slidecraft.agents.builder import BuilderAgent, hex_to_rgb
from slidecraft.config import SlideCraftConfig
from slidecraft.models import (
    ColorPalette,
    DesignSpec,
    DesignStyle,
    FontScheme,
    PresentationContent,
    PresentationPlan,
    SlideContent,
    SlideOutline,
    SlideType,
)


@pytest.fixture
def config():
    return SlideCraftConfig()


@pytest.fixture
def builder(config):
    return BuilderAgent(config)


@pytest.fixture
def design():
    return DesignSpec(
        style=DesignStyle.TECH,
        colors=ColorPalette(
            primary="#6366f1",
            secondary="#06b6d4",
            accent="#f43f5e",
            background="#0f172a",
            text_primary="#f1f5f9",
            text_secondary="#94a3b8",
        ),
        fonts=FontScheme(
            title_font="Arial",
            body_font="Arial",
            title_size_pt=36,
            body_size_pt=18,
        ),
    )


@pytest.fixture
def full_plan():
    """A plan with every slide type."""
    return PresentationPlan(
        title="Test Presentation - All Layouts",
        subtitle="Testing every layout template",
        narrative_arc="Cover → TOC → Sections → Content → End",
        target_audience="Developers",
        slides=[
            SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Cover Slide"),
            SlideOutline(slide_number=2, slide_type=SlideType.TABLE_OF_CONTENTS, title="目录"),
            SlideOutline(slide_number=3, slide_type=SlideType.SECTION_HEADER, title="Part 1"),
            SlideOutline(slide_number=4, slide_type=SlideType.CONTENT, title="Content Page",
                         key_points=["Point A", "Point B"]),
            SlideOutline(slide_number=5, slide_type=SlideType.TWO_COLUMN, title="Comparison"),
            SlideOutline(slide_number=6, slide_type=SlideType.IMAGE_FULL, title="Full Image"),
            SlideOutline(slide_number=7, slide_type=SlideType.IMAGE_TEXT, title="Image + Text"),
            SlideOutline(slide_number=8, slide_type=SlideType.CHART, title="Data Chart"),
            SlideOutline(slide_number=9, slide_type=SlideType.COMPARISON, title="A vs B"),
            SlideOutline(slide_number=10, slide_type=SlideType.TIMELINE, title="Timeline"),
            SlideOutline(slide_number=11, slide_type=SlideType.QUOTE, title="Quote"),
            SlideOutline(slide_number=12, slide_type=SlideType.THANK_YOU, title="Thank You"),
        ],
    )


@pytest.fixture
def full_content():
    """Content matching the full_plan."""
    return PresentationContent(slides=[
        SlideContent(slide_number=1, title="Cover Slide", subtitle="A great subtitle"),
        SlideContent(slide_number=2, title="目录"),
        SlideContent(slide_number=3, title="Part 1: Introduction"),
        SlideContent(
            slide_number=4, title="Key Findings",
            bullet_points=["Finding 1: 50% growth", "Finding 2: New markets", "Finding 3: AI adoption"],
            notes="Emphasize the growth metric",
        ),
        SlideContent(
            slide_number=5, title="Before vs After",
            bullet_points=["Old approach: manual", "New approach: automated", "Result: 3x faster"],
            body_text="Left: Traditional | Right: AI-Powered",
        ),
        SlideContent(slide_number=6, title="Our Vision"),
        SlideContent(
            slide_number=7, title="Product Showcase",
            bullet_points=["Feature A", "Feature B", "Feature C"],
        ),
        SlideContent(
            slide_number=8, title="Growth Data",
            data={"labels": ["Q1", "Q2", "Q3", "Q4"], "values": [10, 25, 45, 80]},
        ),
        SlideContent(
            slide_number=9, title="Plan A vs Plan B",
            bullet_points=["Plan A: Conservative", "Plan A: Low risk", "Plan B: Aggressive", "Plan B: High reward"],
        ),
        SlideContent(
            slide_number=10, title="Our Journey",
            bullet_points=["2023: Founded", "2024: Product Launch", "2025: Global Expansion"],
        ),
        SlideContent(
            slide_number=11, title="Innovation",
            body_text="The best way to predict the future is to create it. — Peter Drucker",
        ),
        SlideContent(slide_number=12, title="Thank You!", subtitle="Questions?"),
    ])


class TestHexToRgb:
    def test_basic(self):
        rgb = hex_to_rgb("#FF0000")
        assert rgb == (0xFF, 0x00, 0x00)

    def test_lowercase(self):
        rgb = hex_to_rgb("#00ff00")
        assert rgb == (0x00, 0xFF, 0x00)

    def test_no_hash(self):
        rgb = hex_to_rgb("0000FF")
        assert rgb == (0x00, 0x00, 0xFF)


class TestBuilderBasic:
    def test_build_simple_presentation(self, builder, design):
        """Test building a minimal presentation with cover + thank you."""
        plan = PresentationPlan(
            title="Simple Test",
            subtitle="Testing basic build",
            slides=[
                SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Hello"),
                SlideOutline(slide_number=2, slide_type=SlideType.THANK_YOU, title="Bye"),
            ],
        )
        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="Hello World", subtitle="Greetings"),
            SlideContent(slide_number=2, title="Thank You!", subtitle="Q&A"),
        ])

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, "test.pptx")
            result = builder.build(plan, content, design, [], out_path)

            assert os.path.exists(result)
            prs = PptxPresentation(result)
            assert len(prs.slides) == 2

    def test_build_all_layouts(self, builder, design, full_plan, full_content):
        """Test building with every slide type."""
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, "all_layouts.pptx")
            result = builder.build(full_plan, full_content, design, [], out_path)

            assert os.path.exists(result)
            prs = PptxPresentation(result)
            assert len(prs.slides) == 12

    def test_output_directory_creation(self, builder, design):
        """Test that output directory is created if it doesn't exist."""
        plan = PresentationPlan(
            title="Dir Test",
            slides=[SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Test")],
        )
        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="Test"),
        ])

        with tempfile.TemporaryDirectory() as tmpdir:
            nested = os.path.join(tmpdir, "a", "b", "c")
            out_path = os.path.join(nested, "test.pptx")
            result = builder.build(plan, content, design, [], out_path)
            assert os.path.exists(result)

    def test_slide_with_notes(self, builder, design):
        """Test that speaker notes are included."""
        plan = PresentationPlan(
            title="Notes Test",
            slides=[
                SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Cover"),
                SlideOutline(slide_number=2, slide_type=SlideType.CONTENT, title="Content"),
            ],
        )
        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="Cover"),
            SlideContent(
                slide_number=2, title="Content with Notes",
                bullet_points=["Point 1"],
                notes="Remember to explain this carefully",
            ),
        ])

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, "notes.pptx")
            result = builder.build(plan, content, design, [], out_path)

            prs = PptxPresentation(result)
            # Content slide should have notes
            content_slide = prs.slides[1]
            notes_text = content_slide.notes_slide.notes_text_frame.text
            assert "explain" in notes_text


class TestBuilderDesignStyles:
    """Test that different design styles produce valid .pptx files."""

    @pytest.mark.parametrize("style", list(DesignStyle))
    def test_all_styles_build(self, builder, style):
        design = DesignSpec(style=style)
        plan = PresentationPlan(
            title=f"Style: {style.value}",
            slides=[
                SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Cover"),
                SlideOutline(slide_number=2, slide_type=SlideType.CONTENT, title="Content"),
            ],
        )
        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="Cover", subtitle="Testing"),
            SlideContent(slide_number=2, title="Content", bullet_points=["Test point"]),
        ])

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, f"style_{style.value}.pptx")
            result = builder.build(plan, content, design, [], out_path)
            assert os.path.exists(result)
            prs = PptxPresentation(result)
            assert len(prs.slides) == 2
