"""
Tests for the enhanced Builder Agent layout templates.
Verifies that each of the 12 slide types renders correctly.
"""

import os
import tempfile
import pytest

from pptx import Presentation as PptxPresentation

from slidecraft.agents.builder import BuilderAgent
from slidecraft.config import SlideCraftConfig
from slidecraft.models import (
    DesignSpec,
    DesignStyle,
    PresentationContent,
    PresentationPlan,
    SlideContent,
    SlideOutline,
    SlideType,
)


@pytest.fixture
def config():
    return SlideCraftConfig()


@pytest.fixture
def builder(config):
    return BuilderAgent(config)


@pytest.fixture
def design():
    return DesignSpec(style=DesignStyle.TECH)


def _build_single_slide(builder, design, slide_type, sc, plan=None):
    """Helper: build a pptx with a single slide of the given type."""
    if plan is None:
        plan = PresentationPlan(
            title="Layout Test",
            slides=[SlideOutline(slide_number=1, slide_type=slide_type, title=sc.title)],
        )
    content = PresentationContent(slides=[sc])
    with tempfile.TemporaryDirectory() as tmpdir:
        out = os.path.join(tmpdir, "test.pptx")
        result = builder.build(plan, content, design, [], out)
        prs = PptxPresentation(result)
        return prs


class TestTwoColumnLayout:
    def test_even_split(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Two Column",
            bullet_points=["Left 1", "Left 2", "Right 1", "Right 2"],
        )
        prs = _build_single_slide(builder, design, SlideType.TWO_COLUMN, sc)
        assert len(prs.slides) == 1
        # Should have multiple shapes (bg elements + text boxes)
        assert len(prs.slides[0].shapes) > 4

    def test_odd_split(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Odd Split",
            bullet_points=["A", "B", "C"],
        )
        prs = _build_single_slide(builder, design, SlideType.TWO_COLUMN, sc)
        assert len(prs.slides) == 1


class TestImageFullLayout:
    def test_without_image(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Full Image", subtitle="A scenic background",
        )
        prs = _build_single_slide(builder, design, SlideType.IMAGE_FULL, sc)
        assert len(prs.slides) == 1

    def test_with_body_text(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Full Image",
            body_text="This is overlay text on the image",
        )
        prs = _build_single_slide(builder, design, SlideType.IMAGE_FULL, sc)
        assert len(prs.slides) == 1


class TestImageTextLayout:
    def test_with_bullets(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Image + Text",
            bullet_points=["Feature A", "Feature B", "Feature C"],
        )
        prs = _build_single_slide(builder, design, SlideType.IMAGE_TEXT, sc)
        assert len(prs.slides) == 1

    def test_with_body_text(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Image + Text",
            body_text="Detailed description of the product.",
        )
        prs = _build_single_slide(builder, design, SlideType.IMAGE_TEXT, sc)
        assert len(prs.slides) == 1


class TestChartLayout:
    def test_with_data(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Revenue Growth",
            data={"labels": ["Q1", "Q2", "Q3", "Q4"], "values": [10, 25, 45, 80]},
        )
        prs = _build_single_slide(builder, design, SlideType.CHART, sc)
        assert len(prs.slides) == 1
        # Verify chart was added (chart shapes exist)
        slide = prs.slides[0]
        has_chart = any(hasattr(s, 'chart') for s in slide.shapes)
        assert has_chart, "Chart shape should be present"

    def test_without_data(self, builder, design):
        """Should fall back to text placeholder if no data."""
        sc = SlideContent(
            slide_number=1, title="Missing Data Chart",
            bullet_points=["See attached data"],
        )
        prs = _build_single_slide(builder, design, SlideType.CHART, sc)
        assert len(prs.slides) == 1

    def test_with_custom_series(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Sales by Region",
            data={
                "labels": ["North", "South", "East", "West"],
                "values": [120, 90, 150, 80],
                "series_name": "销售额(万元)",
            },
        )
        prs = _build_single_slide(builder, design, SlideType.CHART, sc)
        assert len(prs.slides) == 1


class TestComparisonLayout:
    def test_even_comparison(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Plan A vs Plan B",
            bullet_points=[
                "Plan A: Low cost",
                "Plan A: Slower",
                "Plan B: High cost",
                "Plan B: Faster",
            ],
        )
        prs = _build_single_slide(builder, design, SlideType.COMPARISON, sc)
        assert len(prs.slides) == 1
        # Should have card shapes, VS circle, etc.
        assert len(prs.slides[0].shapes) > 5

    def test_single_item(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Only One Point",
            bullet_points=["The only consideration"],
        )
        prs = _build_single_slide(builder, design, SlideType.COMPARISON, sc)
        assert len(prs.slides) == 1


class TestTimelineLayout:
    def test_multiple_milestones(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Company History",
            bullet_points=[
                "2020: Founded",
                "2021: Seed Round",
                "2022: Product Launch",
                "2023: Series A",
                "2024: Global Expansion",
            ],
        )
        prs = _build_single_slide(builder, design, SlideType.TIMELINE, sc)
        assert len(prs.slides) == 1
        # Should have timeline bar + milestone dots + labels
        assert len(prs.slides[0].shapes) > 10

    def test_single_milestone(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Key Moment",
            bullet_points=["2025: The breakthrough"],
        )
        prs = _build_single_slide(builder, design, SlideType.TIMELINE, sc)
        assert len(prs.slides) == 1


class TestQuoteLayout:
    def test_with_attribution(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Quote",
            subtitle="Peter Drucker",
            body_text="The best way to predict the future is to create it.",
        )
        prs = _build_single_slide(builder, design, SlideType.QUOTE, sc)
        assert len(prs.slides) == 1

    def test_without_attribution(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Innovation is key",
        )
        prs = _build_single_slide(builder, design, SlideType.QUOTE, sc)
        assert len(prs.slides) == 1


class TestThankYouLayout:
    def test_with_contact_info(self, builder, design):
        sc = SlideContent(
            slide_number=1, title="Thank You!",
            subtitle="Any Questions?",
            bullet_points=["email@example.com", "www.example.com", "WeChat: exampleID"],
        )
        prs = _build_single_slide(builder, design, SlideType.THANK_YOU, sc)
        assert len(prs.slides) == 1

    def test_minimal(self, builder, design):
        sc = SlideContent(slide_number=1, title="谢谢！")
        prs = _build_single_slide(builder, design, SlideType.THANK_YOU, sc)
        assert len(prs.slides) == 1


class TestFullPresentationAllLayouts:
    """Integration test: build a presentation using every layout type."""

    def test_complete_presentation(self, builder, design):
        plan = PresentationPlan(
            title="Complete Layout Demo",
            subtitle="All 12 slide types",
            slides=[
                SlideOutline(slide_number=1, slide_type=SlideType.COVER, title="Cover"),
                SlideOutline(slide_number=2, slide_type=SlideType.TABLE_OF_CONTENTS, title="TOC"),
                SlideOutline(slide_number=3, slide_type=SlideType.SECTION_HEADER, title="Section"),
                SlideOutline(slide_number=4, slide_type=SlideType.CONTENT, title="Content"),
                SlideOutline(slide_number=5, slide_type=SlideType.TWO_COLUMN, title="Columns"),
                SlideOutline(slide_number=6, slide_type=SlideType.IMAGE_FULL, title="Image"),
                SlideOutline(slide_number=7, slide_type=SlideType.IMAGE_TEXT, title="Img+Txt"),
                SlideOutline(slide_number=8, slide_type=SlideType.CHART, title="Chart"),
                SlideOutline(slide_number=9, slide_type=SlideType.COMPARISON, title="Compare"),
                SlideOutline(slide_number=10, slide_type=SlideType.TIMELINE, title="Timeline"),
                SlideOutline(slide_number=11, slide_type=SlideType.QUOTE, title="Quote"),
                SlideOutline(slide_number=12, slide_type=SlideType.THANK_YOU, title="End"),
            ],
        )
        content = PresentationContent(slides=[
            SlideContent(slide_number=1, title="Welcome", subtitle="Demo Presentation"),
            SlideContent(slide_number=2, title="目录"),
            SlideContent(slide_number=3, title="Part One", subtitle="Introduction"),
            SlideContent(slide_number=4, title="Key Points", bullet_points=["A", "B", "C"]),
            SlideContent(slide_number=5, title="Pros & Cons",
                         bullet_points=["Pro 1", "Pro 2", "Con 1", "Con 2"]),
            SlideContent(slide_number=6, title="Vision"),
            SlideContent(slide_number=7, title="Product", bullet_points=["Feature X"]),
            SlideContent(slide_number=8, title="Growth",
                         data={"labels": ["Jan", "Feb"], "values": [10, 20]}),
            SlideContent(slide_number=9, title="A vs B",
                         bullet_points=["A: Fast", "B: Stable"]),
            SlideContent(slide_number=10, title="Roadmap",
                         bullet_points=["Phase 1", "Phase 2", "Phase 3"]),
            SlideContent(slide_number=11, title="Inspiration",
                         body_text="Think different.", subtitle="Apple"),
            SlideContent(slide_number=12, title="Thank You!", subtitle="Q&A"),
        ])

        with tempfile.TemporaryDirectory() as tmpdir:
            out = os.path.join(tmpdir, "complete.pptx")
            result = builder.build(plan, content, design, [], out)

            assert os.path.exists(result)
            prs = PptxPresentation(result)
            assert len(prs.slides) == 12

            # Verify file size is reasonable (should be > 30KB with all layouts)
            file_size = os.path.getsize(result)
            assert file_size > 30000, f"File too small: {file_size} bytes"
