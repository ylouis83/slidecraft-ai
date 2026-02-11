
"""
Builder Agent — Assembles the final .pptx file using python-pptx with HIGH-FIDELITY styling.

Key Improvements:
- Gradient backgrounds
- Card-based layouts with shadows
- Modern typography handling
- Advanced visual elements (glassmorphism simulation)
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

from slidecraft.agents.base import BaseAgent
from slidecraft.models import (
    DesignSpec,
    GeneratedImage,
    PresentationContent,
    PresentationPlan,
    SlideType,
    DesignStyle,
)

# Import our new advanced styling utilities
# Assuming slidecraft is installed in editable mode or PYTHONPATH includes it
try:
    from slidecraft.utils.pptx_xml import (
        apply_drop_shadow,
        apply_gradient_fill,
        apply_reflection,
        apply_soft_glow,
        apply_text_shadow,
        SubElement
    )
except ImportError:
    # Fallback or local import if running as script
    import sys
    import os
    sys.path.append(os.path.join(os.path.dirname(__file__), "..", ".."))
    from slidecraft.utils.pptx_xml import (
        apply_drop_shadow,
        apply_gradient_fill,
        apply_reflection,
        apply_soft_glow,
        apply_text_shadow,
        SubElement
    )

logger = logging.getLogger("slidecraft")


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_slide_number(slide, prs, design: DesignSpec, number: int):
    """Add a small slide number in the bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.5),
        Inches(0.5), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(10)
    p.font.name = "Arial" # Neutral font for numbers
    p.font.color.rgb = hex_to_rgb(design.colors.text_secondary)
    p.alignment = PP_ALIGN.RIGHT


class BuilderAgent(BaseAgent):
    """PPT builder that assembles the final .pptx file."""

    name = "builder"
    description = "PPT 构建师，将所有素材组装为 .pptx 文件"

    @property
    def system_prompt(self) -> str:
        return "You are the PPT builder agent."

    def build(
        self,
        plan: PresentationPlan,
        content: PresentationContent,
        design: DesignSpec,
        images: list[GeneratedImage],
        output_path: str,
    ) -> str:
        """Build the final .pptx file. Returns the output file path."""
        prs = Presentation()
        prs.slide_width = Inches(design.slide_width_inches)
        prs.slide_height = Inches(design.slide_height_inches)

        # Create lookups
        image_map = {img.slide_number: img for img in images}
        content_map = {sc.slide_number: sc for sc in content.slides}

        # Dispatch table for slide types
        builders = {
            SlideType.COVER: self._build_cover,
            SlideType.THANK_YOU: self._build_thank_you,
            SlideType.SECTION_HEADER: self._build_section,
            SlideType.TABLE_OF_CONTENTS: self._build_toc,
            SlideType.CONTENT: self._build_content,
            SlideType.TWO_COLUMN: self._build_two_column,
            SlideType.IMAGE_FULL: self._build_image_full,
            SlideType.IMAGE_TEXT: self._build_image_text,
            SlideType.CHART: self._build_chart,
            SlideType.COMPARISON: self._build_comparison,
            SlideType.TIMELINE: self._build_timeline,
            SlideType.QUOTE: self._build_quote,
        }

        for i, slide_outline in enumerate(plan.slides):
            sn = slide_outline.slide_number
            sc = content_map.get(sn)
            img = image_map.get(sn)

            builder_fn = builders.get(slide_outline.slide_type, self._build_content)

            # Each builder has a slightly different signature; use kwargs
            slide = None
            if slide_outline.slide_type == SlideType.COVER:
                slide = builder_fn(prs, plan, sc, design, img)
            elif slide_outline.slide_type == SlideType.TABLE_OF_CONTENTS:
                slide = builder_fn(prs, plan, sc, design)
            else:
                slide = builder_fn(prs, sc, design, img)
            
            # Add slide number to all except Cover
            if slide and slide_outline.slide_type != SlideType.COVER:
                _add_slide_number(slide, prs, design, sn)

        # Save
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        logger.info(f"[{self.name}] Saved presentation to {out}")
        return str(out)

    # ══════════════════════════════════════════════════════════
    #  Shared helpers
    # ══════════════════════════════════════════════════════════

    def _add_bg(self, slide, design: DesignSpec, prs=None):
        """Add background color or gradient to slide."""
        if design.use_gradient_backgrounds and prs:
            # Use a full-slide rectangle for gradient
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = hex_to_rgb(design.colors.background)
            bg_shape.line.fill.background() # No border
            
            # Apply Gradient
            apply_gradient_fill(
                bg_shape, 
                design.colors.background, 
                design.colors.surface, 
                angle=45
            )
        else:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(design.colors.background)

    def _add_top_bar(self, slide, prs, design: DesignSpec, height: float = 0.08):
        """Add a neon/gradient accent bar at top of slide."""
        # Main bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
        bar.line.fill.background()
        
        # Glow effect
        if design.style in [DesignStyle.DARK, DesignStyle.TECH]:
            apply_soft_glow(bar, design.colors.primary, alpha="60000")

    def _add_title(self, slide, title: str, design: DesignSpec,
                   left=0.8, top=0.5, width=11.0, height=0.8):
        """Add a standard title textbox with proper font execution."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = design.fonts.title_font
        p.font.size = Pt(design.fonts.title_size_pt)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(design.colors.text_primary)
        
        # Add subtle shadow to title in dark mode for pop
        if design.style == DesignStyle.DARK:
             pass # python-pptx text shadow is tricky, handled by shape usually

        return txBox

    def _add_card_background(self, slide, left, top, width, height, design: DesignSpec):
        """Add a rounded rect card background with shadow."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(design.colors.surface)
        card.line.fill.background() # No border
        
        # Adjust rounded corners (default is too round)
        card.adjustments[0] = 0.05
        
        # Add Drop Shadow
        apply_drop_shadow(card, transparency=70, blur=15, distance=5)
        
        return card

    def _add_bullets(self, slide, bullets: list[str], design: DesignSpec,
                     left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), height=Inches(5.0)):
        """Add a bullet-point text frame. Args must be in EMUs."""
        txBox = slide.shapes.add_textbox(
            left, top, width, height
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bp in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bp}"
            p.font.name = design.fonts.body_font
            p.font.size = Pt(design.fonts.body_size_pt)
            p.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            p.space_after = Pt(14)
            p.line_spacing = 1.2 # Better readability
        return txBox

    def _add_notes(self, slide, notes: str):
        """Add speaker notes to a slide."""
        if notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes

    def _try_add_image(self, slide, img: Optional[GeneratedImage],
                       left=Inches(8.5), top=Inches(1.5), width=Inches(4.3), height=Inches(4.3)):
        """Try to add an image to the slide, ignore if file missing. 
        Args must be in EMUs (use Inches() at call site)."""
        if img and Path(img.image_path).exists():
            try:
                pic = slide.shapes.add_picture(
                    img.image_path,
                    left, top,
                    width=width, height=height,
                )
                # Add shadow and border to image
                line = pic.line
                line.color.rgb = RGBColor(100, 100, 100)
                line.width = Pt(1)
                apply_drop_shadow(pic, transparency=50, blur=10, distance=5)
                return True
            except Exception as e:
                logger.warning(f"Could not add image: {e}")
        return False

    # ══════════════════════════════════════════════════════════
    #  1. COVER
    # ══════════════════════════════════════════════════════════

    def _build_cover(self, prs, plan, sc, design, img=None):
        """Cover slide with large title, subtitle, and optional accent."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        
        # 1. Background Image (if available)
        if img and Path(img.image_path).exists():
            pic = slide.shapes.add_picture(img.image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            # Add Dark Overlay
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = hex_to_rgb("#000000")
            overlay.fill.transparency = 0.3 # 30% transparent (70% opacity in PPT terms? No, 0.0 is opaque)
            # wait, python-pptx transparency is ... tricky. 
            # fill.transparency is float 0.0 to 1.0 (0 is opaque)
            # But XML might be different. Let's use solid color for now or rely on gradient if no image.
            # Actually, let's use a semi-transparent gradient or shape.
            overlay.line.fill.background()
        else:
            self._add_bg(slide, design, prs)

        w = prs.slide_width
        h = prs.slide_height

        # Dynamic Gradient Background Layer (only if no image)
        if not (img and Path(img.image_path).exists()):
             # Geometric Shapes for "Tech" feel
            shape_bg = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, 
                w - Inches(6), -Inches(2), Inches(8), Inches(10)
            )
            shape_bg.rotation = 270
            shape_bg.fill.solid()
            shape_bg.fill.fore_color.rgb = hex_to_rgb(design.colors.surface)
            shape_bg.line.fill.background()
            # Add gradient to shape
            apply_gradient_fill(shape_bg, design.colors.surface, design.colors.background, angle=90)

        # Title Area
        title = sc.title if sc else plan.title
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.5), Inches(9), Inches(2.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = design.fonts.title_font
        p.font.size = Pt(54) # Larger title
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(design.colors.primary)
        p.alignment = PP_ALIGN.LEFT
        
        # Subtitle
        subtitle = sc.subtitle if sc else plan.subtitle
        if subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.2), Inches(9), Inches(1)
            )
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.name = design.fonts.body_font
            p2.font.size = Pt(28)
            p2.font.color.rgb = hex_to_rgb(design.colors.text_secondary)

        # Footer / Organization
        txBox3 = slide.shapes.add_textbox(
            Inches(1.0), h - Inches(1.0), Inches(5), Inches(0.5)
        )
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "CONFIDENTIAL | 2026"
        p3.font.size = Pt(12)
        p3.font.color.rgb = hex_to_rgb(design.colors.text_secondary)
        
        # Logo placeholder (circle) top right
        logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, w - Inches(1.5), Inches(0.5), Inches(1), Inches(1))
        logo.fill.solid()
        logo.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
        apply_soft_glow(logo, design.colors.primary)

        return slide

    # ══════════════════════════════════════════════════════════
    #  2. TABLE OF CONTENTS
    # ══════════════════════════════════════════════════════════

    def _build_toc(self, prs, plan, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        self._add_top_bar(slide, prs, design)

        # Title
        self._add_title(slide, sc.title if sc else "目录", design, width=5.0)

        # 2-Column Grid Layout for TOC
        sections = [
            s for s in plan.slides
            if s.slide_type in (SlideType.SECTION_HEADER, SlideType.CONTENT)
        ]
        
        start_x = Inches(1.0)
        start_y = Inches(1.8)
        col_width = (prs.slide_width - Inches(2.0)) / 2
        
        for i, sec in enumerate(sections[:8], 1):
            col = (i-1) // 4
            row = (i-1) % 4
            
            x = start_x + (col * col_width)
            y = start_y + (row * Inches(1.2))
            
            # Card for each item
            card = self._add_card_background(slide, x, y, col_width - Inches(0.5), Inches(1.0), design)
            
            # Number
            num_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.8))
            np = num_box.text_frame.paragraphs[0]
            np.text = f"{i:02d}"
            np.font.size = Pt(32)
            np.font.bold = True
            np.font.color.rgb = hex_to_rgb(design.colors.accent) # Accent color for number
            
            # Text
            txt_box = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.25), col_width - Inches(1.5), Inches(0.6))
            tp = txt_box.text_frame.paragraphs[0]
            tp.text = sec.title
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            
        return slide

    # ══════════════════════════════════════════════════════════
    #  3. SECTION HEADER
    # ══════════════════════════════════════════════════════════

    def _build_section(self, prs, sc, design, img=None):
        """Section header with full-width dynamic gradient band."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        if img and Path(img.image_path).exists():
           pic = slide.shapes.add_picture(img.image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
           # Dark overlay
           overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
           overlay.fill.solid()
           overlay.fill.fore_color.rgb = hex_to_rgb("#000000")
           overlay.fill.transparency = 0.5
           overlay.line.fill.background()
        else:
           self._add_bg(slide, design, prs)

        w = prs.slide_width
        h = prs.slide_height

        # Diagonal large shape (only if no image)
        if not (img and Path(img.image_path).exists()):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE, 
                0, 0, w, h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(design.colors.surface)
            # Gradient
            apply_gradient_fill(shape, design.colors.background, design.colors.surface, angle=45)

        # Central Box
        box_h = Inches(3.0)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, (h - box_h)/2, w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
        box.line.fill.background()
        apply_drop_shadow(box)
        
        # Section Number/Subtitle
        if sc and sc.subtitle:
            st_box = slide.shapes.add_textbox(Inches(1.0), (h - box_h)/2 + Inches(0.5), w - Inches(2), Inches(0.5))
            stp = st_box.text_frame.paragraphs[0]
            stp.text = sc.subtitle
            stp.font.size = Pt(18)
            stp.font.color.rgb = hex_to_rgb(design.colors.accent)
            stp.alignment = PP_ALIGN.LEFT

        # Title
        if sc:
            t_box = slide.shapes.add_textbox(Inches(1.0), (h - box_h)/2 + Inches(1.0), w - Inches(2), Inches(1.5))
            tp = t_box.text_frame.paragraphs[0]
            tp.text = sc.title
            tp.font.name = design.fonts.title_font
            tp.font.size = Pt(48)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(255, 255, 255) # Always white on primary
            apply_text_shadow(t_box.text_frame.paragraphs[0].runs[0])
            
        return slide

    # ══════════════════════════════════════════════════════════
    #  4. CONTENT (standard)
    # ══════════════════════════════════════════════════════════

    def _build_content(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        if not sc: return slide

        self._add_top_bar(slide, prs, design)
        self._add_title(slide, sc.title, design)

        # Layout Logic: Card
        has_image = img and Path(img.image_path).exists() if img else False
        
        content_left = Inches(0.8)
        content_top = Inches(1.5)
        content_w = Inches(11.7) if not has_image else Inches(7.0)
        content_h = Inches(5.2)

        # Card Background
        self._add_card_background(slide, content_left, content_top, content_w, content_h, design)

        # Bullets
        if sc.bullet_points:
            self._add_bullets(
                slide, sc.bullet_points, design,
                left=content_left + Inches(0.3), 
                top=content_top + Inches(0.3), 
                width=content_w - Inches(0.6), 
                height=content_h - Inches(0.6)
            )
        elif sc.body_text:
            # Paragraph text
            tb = slide.shapes.add_textbox(content_left + Inches(0.3), content_top + Inches(0.3), content_w - Inches(0.6), content_h - Inches(0.6))
            tp = tb.text_frame.paragraphs[0]
            tp.text = sc.body_text
            tp.font.size = Pt(design.fonts.body_size_pt)
            tp.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            tp.word_wrap = True

        # Image
        if has_image:
            # Calculate remaining width in inches first
            remaining_w_inches = 12.5 - (content_left.inches + content_w.inches) - 0.2
            self._try_add_image(
                slide, img, 
                left=content_left + content_w + Inches(0.2), 
                top=content_top, 
                width=Inches(remaining_w_inches), 
                height=content_h
            )

        self._add_notes(slide, sc.notes)
        return slide

    # ══════════════════════════════════════════════════════════
    #  5. TWO COLUMN
    # ══════════════════════════════════════════════════════════

    def _build_two_column(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background Image (if available) with Overlay
        if img and Path(img.image_path).exists():
           pic = slide.shapes.add_picture(img.image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
           overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
           overlay.fill.solid()
           overlay.fill.fore_color.rgb = hex_to_rgb("#000000")
           overlay.fill.transparency = 0.5
           overlay.line.fill.background()
        else:
           self._add_bg(slide, design, prs)
           
        if not sc: return slide

        self._add_top_bar(slide, prs, design)
        self._add_title(slide, sc.title, design)

        points = sc.bullet_points or []
        mid = (len(points) + 1) // 2
        left_points = points[:mid]
        right_points = points[mid:]

        # Two Cards
        col_w = Inches(5.7)
        col_h = Inches(5.0)
        gap = Inches(0.3)
        start_x = Inches(0.8)
        start_y = Inches(1.6)

        # Left Card
        self._add_card_background(slide, start_x, start_y, col_w, col_h, design)
        if left_points:
            self._add_bullets(slide, left_points, design, start_x + Inches(0.2), start_y + Inches(0.2), col_w - Inches(0.4), col_h - Inches(0.4))

        # Right Card
        self._add_card_background(slide, start_x + col_w + gap, start_y, col_w, col_h, design)
        if right_points:
            self._add_bullets(slide, right_points, design, start_x + col_w + gap + Inches(0.2), start_y + Inches(0.2), col_w - Inches(0.4), col_h - Inches(0.4))

        self._add_notes(slide, sc.notes)
        return slide

    # ══════════════════════════════════════════════════════════
    #  6. IMAGE FULL
    # ══════════════════════════════════════════════════════════

    def _build_image_full(self, prs, sc, design, img=None):
        """Full bleed image with glass overlay."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        w, h = prs.slide_width, prs.slide_height

        if img and Path(img.image_path).exists():
            pic = slide.shapes.add_picture(img.image_path, 0, 0, width=w, height=h)
        else:
            # Fallback bg
            self._add_bg(slide, design, prs)

        # Glass Overlay at bottom
        overlay_h = Inches(2.5)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, h - overlay_h, w, overlay_h)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = hex_to_rgb(design.colors.background)
        # Simulate transparency via alpha in XML? 
        # For simplicity, using solid dark color
        
        # Text
        if sc:
            tb = slide.shapes.add_textbox(Inches(1), h - overlay_h + Inches(0.5), w - Inches(2), Inches(1))
            tp = tb.text_frame.paragraphs[0]
            tp.text = sc.title
            tp.font.size = Pt(36)
            tp.font.bold = True
            tp.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            
            if sc.body_text or sc.subtitle:
                tb2 = slide.shapes.add_textbox(Inches(1), h - overlay_h + Inches(1.2), w - Inches(2), Inches(1))
                tp2 = tb2.text_frame.paragraphs[0]
                tp2.text = sc.body_text or sc.subtitle
                tp2.font.size = Pt(18)
                tp2.font.color.rgb = hex_to_rgb(design.colors.text_secondary)

        return slide

    # ══════════════════════════════════════════════════════════
    #  7. IMAGE TEXT
    # ══════════════════════════════════════════════════════════

    def _build_image_text(self, prs, sc, design, img=None):
        return self._build_content(prs, sc, design, img) # Reuse content builder which handles images nicely now

    # ══════════════════════════════════════════════════════════
    #  8. CHART
    # ══════════════════════════════════════════════════════════

    def _build_chart(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        if not sc: return slide
        
        self._add_top_bar(slide, prs, design)
        self._add_title(slide, sc.title, design)
        
        # Chart Container Card
        card_rect = self._add_card_background(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), design)

        if sc.data and "labels" in sc.data:
            from pptx.chart.data import CategoryChartData
            chart_data = CategoryChartData()
            chart_data.categories = sc.data["labels"]
            chart_data.add_series(sc.data.get("series_name", "Series"), sc.data["values"])
            
            x, y, cx, cy = Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.6)
            graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
            chart = graphic_frame.chart
            
            # Simple styling
            try:
                chart.plots[0].series[0].format.fill.solid()
                chart.plots[0].series[0].format.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
            except:
                pass
        else:
            # Fallback
            tb = slide.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(1))
            tb.text_frame.text = "No Chart Data"
            
        return slide

    # ══════════════════════════════════════════════════════════
    #  9. COMPARISON
    # ══════════════════════════════════════════════════════════

    def _build_comparison(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        if not sc: return slide
        
        self._add_top_bar(slide, prs, design)
        self._add_title(slide, sc.title, design)
        
        points = sc.bullet_points or []
        mid = max(len(points) // 2, 1)
        left_items = points[:mid]
        right_items = points[mid:]
        
        # Left Card
        self._add_card_background(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), design)
        # Header A
        bar_a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.6))
        bar_a.fill.solid()
        bar_a.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
        bar_a.text_frame.text = "Before / Initial"
        bar_a.text_frame.paragraphs[0].font.bold = True
        
        if left_items:
            self._add_bullets(slide, left_items, design, Inches(1.0), Inches(2.6), Inches(5.2), Inches(3.8))

        # Right Card
        self._add_card_background(slide, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), design)
        # Header B
        bar_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.6))
        bar_b.fill.solid()
        bar_b.fill.fore_color.rgb = hex_to_rgb(design.colors.secondary)
        bar_b.text_frame.text = "After / Future"
        bar_b.text_frame.paragraphs[0].font.bold = True

        if right_items:
            self._add_bullets(slide, right_items, design, Inches(7.1), Inches(2.6), Inches(5.2), Inches(3.8))
        
        return slide

    # ══════════════════════════════════════════════════════════
    #  10. TIMELINE
    # ══════════════════════════════════════════════════════════

    def _build_timeline(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        if not sc: return slide
        
        self._add_top_bar(slide, prs, design)
        self._add_title(slide, sc.title, design)
        
        points = sc.bullet_points or []
        
        # Central Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(11.3), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(design.colors.primary)
        apply_soft_glow(line, design.colors.primary)
        
        count = len(points)
        step = Inches(11.3) / max(count, 1)
        
        for i, pt_text in enumerate(points):
            cx = Inches(1) + step * i + step/2
            cy = Inches(3.55)
            
            # Dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.2), cy - Inches(0.2), Inches(0.4), Inches(0.4))
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb(design.colors.accent)
            dot.line.fill.background()
            
            # Text Box
            is_top = i % 2 == 0
            ty = cy - Inches(1.5) if is_top else cy + Inches(0.5)
            
            tb = slide.shapes.add_textbox(cx - Inches(1.5), ty, Inches(3), Inches(1))
            tp = tb.text_frame.paragraphs[0]
            tp.text = pt_text
            tp.alignment = PP_ALIGN.CENTER
            tp.font.size = Pt(12)
            tp.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            
        return slide

    # ══════════════════════════════════════════════════════════
    #  11. QUOTE
    # ══════════════════════════════════════════════════════════

    def _build_quote(self, prs, sc, design, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        
        # Giant Quote Mark
        q = slide.shapes.add_shape(MSO_SHAPE.DOUBLE_WAVE, Inches(1), Inches(1), Inches(2), Inches(2)) # Abstract shape
        q.fill.solid()
        q.fill.fore_color.rgb = hex_to_rgb(design.colors.accent)
        
        if sc:
            tb = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.3), Inches(3))
            tp = tb.text_frame.paragraphs[0]
            tp.text = f'"{sc.body_text}"'
            tp.font.size = Pt(32)
            tp.font.italic = True
            tp.font.name = design.fonts.title_font
            tp.font.color.rgb = hex_to_rgb(design.colors.text_primary)
            tp.alignment = PP_ALIGN.CENTER
            
            if sc.subtitle:
                tb2 = slide.shapes.add_textbox(Inches(6), Inches(5.5), Inches(5), Inches(1))
                tp2 = tb2.text_frame.paragraphs[0]
                tp2.text = f"— {sc.subtitle}"
                tp2.alignment = PP_ALIGN.RIGHT
                tp2.font.color.rgb = hex_to_rgb(design.colors.text_secondary)

        return slide

    # ══════════════════════════════════════════════════════════
    #  12. THANK YOU
    # ══════════════════════════════════════════════════════════

    def _build_thank_you(self, prs, sc, design, img=None):
        """Centered thank you slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, design, prs)
        
        w, h = prs.slide_width, prs.slide_height
        
        tb = slide.shapes.add_textbox(0, h/2 - Inches(1), w, Inches(2))
        tp = tb.text_frame.paragraphs[0]
        tp.text = sc.title if sc else "Thank You"
        tp.font.size = Pt(60)
        tp.font.bold = True
        tp.font.color.rgb = hex_to_rgb(design.colors.primary)
        tp.alignment = PP_ALIGN.CENTER
        apply_text_shadow(tp.runs[0])
        apply_reflection(tb)
        
        return slide
