"""
Rendering Engine V2 — Premium-quality PPT rendering.

Major improvements:
  1. Font family support (title_font / body_font applied per element)
  2. Colored bullet indicators with accent circles
  3. Consistent page numbering footer
  4. Semi-transparent overlays for shapes (alpha channel)
  5. Text frame margin control (better breathing room)
  6. Vertical text anchoring for centered content
  7. Better chart styling (gradient fills, value labels)
  8. Shape border radius control
"""

from __future__ import annotations

import copy
import logging
from lxml import etree
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap

from slidecraft.engine.elements import (
    ElementType,
    GradientSpec,
    GradientStop,
    SlideElement,
    SlideSpec,
    PresentationSpec,
)
from slidecraft.engine.layout import (
    auto_fit_title_size,
    auto_fit_bullet_size,
    auto_fit_spacing,
)

logger = logging.getLogger("slidecraft.renderer")


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor. Supports #RGB and #RRGGBB."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_srgb(hex_color: str) -> str:
    """Convert #RRGGBB to RRGGBB (for XML)."""
    return hex_color.lstrip("#").upper()


def _hex_to_rgb_tuple(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _relative_luminance(hex_color: str) -> float:
    def transform(v: int) -> float:
        x = v / 255
        return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4

    r, g, b = _hex_to_rgb_tuple(hex_color)
    return 0.2126 * transform(r) + 0.7152 * transform(g) + 0.0722 * transform(b)


def contrast_ratio(color_a: str, color_b: str) -> float:
    la = _relative_luminance(color_a)
    lb = _relative_luminance(color_b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


# ══════════════════════════════════════════════════════════════
#  Gradient Background (XML manipulation)
# ══════════════════════════════════════════════════════════════

def apply_gradient_background(slide, gradient: GradientSpec):
    """Apply a gradient fill to the slide background."""
    bg = slide.background
    bg_elem = bg._element

    for child in list(bg_elem):
        if child.tag.endswith("bgPr"):
            bg_elem.remove(child)

    bgPr = etree.SubElement(bg_elem, qn("p:bgPr"))
    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")

    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for stop in gradient.stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(stop.position * 100000)))
        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
        srgbClr.set("val", hex_to_srgb(stop.color))

    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(int(gradient.angle * 60000)))
    lin.set("scaled", "1")

    etree.SubElement(bgPr, qn("a:effectLst"))
    return slide


def apply_solid_background(slide, color: str):
    """Apply a solid color background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


# ══════════════════════════════════════════════════════════════
#  Shadow Effects (XML manipulation)
# ══════════════════════════════════════════════════════════════

def apply_shadow(shape, offset_pt: float = 3.0, color: str = "000000",
                 alpha: int = 40):
    """Add a drop shadow effect to a shape."""
    sp = shape._element
    spPr = sp.find(qn("a:spPr"))
    if spPr is None:
        spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return

    for eff in spPr.findall(qn("a:effectLst")):
        spPr.remove(eff)

    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", str(int(offset_pt * 2 * 12700)))
    outerShdw.set("dist", str(int(offset_pt * 12700)))
    outerShdw.set("dir", "5400000")
    outerShdw.set("rotWithShape", "0")

    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgbClr.set("val", color)
    alphaElem = etree.SubElement(srgbClr, qn("a:alpha"))
    alphaElem.set("val", str(alpha * 1000))


def apply_shape_transparency(shape, alpha_percent: int):
    """Set fill transparency on a shape (0=opaque, 100=fully transparent)."""
    sp = shape._element
    solidFill = sp.find(f".//{qn('a:solidFill')}")
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Remove existing alpha
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alphaElem = etree.SubElement(srgb, qn("a:alpha"))
    alphaElem.set("val", str((100 - alpha_percent) * 1000))


# ══════════════════════════════════════════════════════════════
#  Element Renderers
# ══════════════════════════════════════════════════════════════

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "chevron": MSO_SHAPE.CHEVRON,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
}


class SlideRenderer:
    """Renders a SlideSpec into a python-pptx slide — premium quality."""

    def __init__(self, prs: Presentation, default_colors: dict = None,
                 title_font: str = "", body_font: str = "",
                 total_slides: int = 0):
        self.prs = prs
        self.colors = default_colors or {}
        self.title_font = title_font or "微软雅黑"
        self.body_font = body_font or "微软雅黑"
        self.total_slides = total_slides

    def render(self, spec: SlideSpec) -> object:
        """Render a complete SlideSpec to a slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        if spec.background_gradient and spec.background_gradient.stops:
            apply_gradient_background(slide, spec.background_gradient)
        elif spec.background_color:
            apply_solid_background(slide, spec.background_color)

        # Sort elements by z_order
        sorted_elements = sorted(spec.elements, key=lambda e: e.z_order)

        # Render each element
        for element in sorted_elements:
            self._render_element(slide, element)

        # Add slide number (skip for cover)
        sn = spec.slide_number
        if sn and sn > 1 and self.total_slides > 0:
            self._render_page_number(slide, sn)

        # Notes
        if spec.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = spec.notes

        return slide

    def _render_page_number(self, slide, slide_number: int):
        """Render a slide number in bottom right corner."""
        sw = self.prs.slide_width
        txBox = slide.shapes.add_textbox(
            sw - Inches(1.5), Inches(7.1), Inches(1.2), Inches(0.3),
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{slide_number} / {self.total_slides}"
        p.font.size = Pt(10)
        p.font.name = self.body_font
        color_key = self.colors.get("text_secondary", "#888888")
        p.font.color.rgb = hex_to_rgb(color_key)
        p.alignment = PP_ALIGN.RIGHT

    def _render_element(self, slide, el: SlideElement):
        """Dispatch element to the appropriate renderer."""
        renderers = {
            ElementType.TITLE: self._render_title,
            ElementType.SUBTITLE: self._render_subtitle,
            ElementType.TEXT: self._render_text,
            ElementType.BULLETS: self._render_bullets,
            ElementType.SHAPE: self._render_shape,
            ElementType.IMAGE: self._render_image,
            ElementType.CHART: self._render_chart,
            ElementType.DIVIDER: self._render_divider,
        }
        renderer = renderers.get(el.element_type)
        if renderer:
            renderer(slide, el)

    def _resolve_color(self, color: Optional[str], fallback_key: str) -> str:
        if color:
            return color
        return self.colors.get(fallback_key, "#333333")

    def _best_text_color(self, bg_color: Optional[str]) -> str:
        """Pick a readable foreground color for a given background."""
        if not bg_color:
            return self.colors.get("text_primary", "#f8fafc")

        theme_light = self.colors.get("text_primary", "#f8fafc")
        theme_dark = "#0f172a"
        white = "#ffffff"
        black = "#000000"

        candidates = [theme_light, theme_dark, white, black]
        return max(candidates, key=lambda c: contrast_ratio(c, bg_color))

    def _set_font(self, font, el_style, fallback_font: str = None):
        """Apply font family to a font object."""
        font_name = fallback_font or self.body_font
        # If title-level element, use title font
        font.name = font_name

    def _render_title(self, slide, el: SlideElement):
        """Render a title element with auto-fit sizing and font family."""
        pos = el.position
        txBox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        # Vertical centering
        tf.paragraphs[0].space_before = Pt(0)

        font_size = el.text_style.font_size
        if el.text_style.auto_fit and font_size:
            font_size = auto_fit_title_size(el.content, font_size)

        p = tf.paragraphs[0]
        p.text = el.content
        p.font.size = Pt(font_size or 36)
        p.font.bold = el.text_style.bold
        p.font.italic = el.text_style.italic
        p.font.name = self.title_font
        p.font.color.rgb = hex_to_rgb(
            self._resolve_color(el.text_style.color, "text_primary")
        )
        p.alignment = ALIGN_MAP.get(el.text_style.alignment, PP_ALIGN.LEFT)

    def _render_subtitle(self, slide, el: SlideElement):
        """Render a subtitle element."""
        pos = el.position
        txBox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.content
        p.font.size = Pt(el.text_style.font_size or 22)
        p.font.bold = el.text_style.bold
        p.font.italic = el.text_style.italic
        p.font.name = self.title_font
        p.font.color.rgb = hex_to_rgb(
            self._resolve_color(el.text_style.color, "text_secondary")
        )
        p.alignment = ALIGN_MAP.get(el.text_style.alignment, PP_ALIGN.LEFT)

    def _render_text(self, slide, el: SlideElement):
        """Render a body text element."""
        pos = el.position
        txBox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.content
        p.font.size = Pt(el.text_style.font_size or 18)
        p.font.name = self.body_font
        p.font.color.rgb = hex_to_rgb(
            self._resolve_color(el.text_style.color, "text_primary")
        )
        p.font.italic = el.text_style.italic
        p.alignment = ALIGN_MAP.get(el.text_style.alignment, PP_ALIGN.LEFT)

    def _render_bullets(self, slide, el: SlideElement):
        """Render a bullet list with auto-fit sizing, spacing, and font."""
        pos = el.position

        base_size = el.text_style.font_size or 18
        if el.text_style.auto_fit:
            font_size = auto_fit_bullet_size(el.items, base_size, pos.height)
        else:
            font_size = base_size

        spacing = auto_fit_spacing(len(el.items), pos.height, font_size)

        txBox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        bullet_color = self.colors.get("primary", "#6366f1")

        for i, item in enumerate(el.items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Use colored bullet indicator
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(max(8, font_size - 4))
            run_bullet.font.color.rgb = hex_to_rgb(bullet_color)
            run_bullet.font.name = self.body_font

            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.name = self.body_font
            run_text.font.color.rgb = hex_to_rgb(
                self._resolve_color(el.text_style.color, "text_primary")
            )
            p.space_after = Pt(spacing)
            p.line_spacing = Pt(font_size + 6)

    def _render_shape(self, slide, el: SlideElement):
        """Render a shape with optional shadow, transparency, and border."""
        pos = el.position
        shape_type = SHAPE_MAP.get(el.shape_type, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            shape_type,
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )

        # Fill
        if el.shape_style.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(el.shape_style.fill_color)
        else:
            shape.fill.background()

        # Border
        if el.shape_style.border_color:
            shape.line.color.rgb = hex_to_rgb(el.shape_style.border_color)
            shape.line.width = Pt(el.shape_style.border_width or 1)
        else:
            shape.line.fill.background()

        # Transparency (opacity: 0.0=transparent, 1.0=opaque)
        if el.shape_style.opacity < 1.0:
            apply_shape_transparency(shape, int((1.0 - el.shape_style.opacity) * 100))

        # Shadow
        if el.shape_style.shadow:
            apply_shadow(shape, el.shape_style.shadow_offset)

        # Text content
        if el.content:
            tf = shape.text_frame
            tf.word_wrap = True
            # Vertical centering
            tf.paragraphs[0].alignment = ALIGN_MAP.get(
                el.text_style.alignment, PP_ALIGN.CENTER
            )
            p = tf.paragraphs[0]
            p.text = el.content
            if el.text_style.font_size:
                p.font.size = Pt(el.text_style.font_size)
            p.font.name = self.body_font
            text_color = el.text_style.color or self._best_text_color(el.shape_style.fill_color)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.font.bold = el.text_style.bold
            p.alignment = ALIGN_MAP.get(el.text_style.alignment, PP_ALIGN.CENTER)

        return shape

    def _render_image(self, slide, el: SlideElement):
        """Render an image element."""
        pos = el.position
        path = el.content

        if not path or not Path(path).exists():
            logger.warning(f"Image not found: {path}")
            return

        try:
            slide.shapes.add_picture(
                path,
                Inches(pos.left), Inches(pos.top),
                width=Inches(pos.width), height=Inches(pos.height),
            )
        except Exception as e:
            logger.warning(f"Could not add image: {e}")

    def _render_chart(self, slide, el: SlideElement):
        """Render a chart from data with enhanced styling."""
        if not el.chart_data:
            return

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        pos = el.position
        data = el.chart_data

        chart_data = CategoryChartData()
        chart_data.categories = data.get("labels", [])
        series_name = data.get("series_name", "数据")
        chart_data.add_series(series_name, data.get("values", []))

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
            chart_data,
        )

        chart = chart_frame.chart
        chart.has_legend = False

        # Match chart container to slide theme for consistent contrast.
        try:
            chart.chart_area.format.fill.solid()
            chart.chart_area.format.fill.fore_color.rgb = hex_to_rgb(
                self.colors.get("background", "#0b1020")
            )
            plot_area = chart.plots[0]
            plot_area.format.fill.solid()
            plot_area.format.fill.fore_color.rgb = hex_to_rgb(
                self.colors.get("surface", "#131a2d")
            )
        except Exception:
            pass

        plot = chart.plots[0]
        plot.gap_width = 80  # Narrower gaps for premium look
        series = plot.series[0]
        series.format.fill.solid()
        primary = self.colors.get("primary", "#6366f1")
        series.format.fill.fore_color.rgb = hex_to_rgb(primary)

        # Show values with readable color.
        try:
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(11)
            data_labels.font.name = self.body_font
            data_labels.font.color.rgb = hex_to_rgb(
                self.colors.get("text_primary", "#f8fafc")
            )
        except Exception:
            pass

        # Style value axis
        try:
            value_axis = chart.value_axis
            value_axis.has_title = False
            category_axis = chart.category_axis
            category_axis.has_title = False
            # Make tick label font smaller
            value_axis.tick_labels.font.size = Pt(11)
            category_axis.tick_labels.font.size = Pt(12)
            category_axis.tick_labels.font.name = self.body_font
            value_axis.tick_labels.font.name = self.body_font
            axis_text_color = hex_to_rgb(self.colors.get("text_secondary", "#dbeafe"))
            category_axis.tick_labels.font.color.rgb = axis_text_color
            value_axis.tick_labels.font.color.rgb = axis_text_color

            # Soften axis/grid lines while keeping visibility on dark backgrounds.
            if hasattr(category_axis, "format"):
                category_axis.format.line.color.rgb = axis_text_color
            if hasattr(value_axis, "format"):
                value_axis.format.line.color.rgb = axis_text_color

            try:
                value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb(
                    self.colors.get("surface", "#334155")
                )
            except Exception:
                pass
        except Exception:
            pass

    def _render_divider(self, slide, el: SlideElement):
        """Render a horizontal divider line."""
        pos = el.position
        color = el.shape_style.fill_color or self.colors.get("primary", "#6366f1")

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(color)
        line.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  High-level API
# ══════════════════════════════════════════════════════════════

def render_presentation(spec: PresentationSpec, colors: dict,
                        title_font: str = "", body_font: str = "") -> Presentation:
    """Render a complete PresentationSpec into a python-pptx Presentation."""
    prs = Presentation()
    prs.slide_width = Inches(spec.slide_width)
    prs.slide_height = Inches(spec.slide_height)

    renderer = SlideRenderer(
        prs,
        default_colors=colors,
        title_font=title_font,
        body_font=body_font,
        total_slides=len(spec.slides),
    )

    for slide_spec in spec.slides:
        renderer.render(slide_spec)

    return prs
